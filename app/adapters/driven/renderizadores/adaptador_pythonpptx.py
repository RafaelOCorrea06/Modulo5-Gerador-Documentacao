# Adaptador driven: renderiza Apresentacao para .pptx via python-pptx.
# Aplica identidade visual Mackenzie (cores configuraveis via settings).

import base64
import io
from typing import Callable, Dict

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from app.application.ports.driven.renderizador_pptx import RenderizadorPPTX
from app.config import settings
from app.domain.entidades.apresentacao import Apresentacao
from app.domain.entidades.slide import Slide, TipoSlide
from app.domain.excecoes import ApresentacaoInvalidaError, RenderizadorError


_LARGURA_SLIDE = Inches(13.333)  # 16:9
_ALTURA_SLIDE = Inches(7.5)
_LAYOUT_BRANCO = 6  # layout em branco padrao do python-pptx

_CORES_INDICADOR = {
    "verde": RGBColor(0x2E, 0x7D, 0x32),
    "amarelo": RGBColor(0xF9, 0xA8, 0x25),
    "vermelho": RGBColor(0xC6, 0x28, 0x28),
}


def _hex_para_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class AdaptadorPythonPPTX(RenderizadorPPTX):
    """Renderizador PPTX com layouts pre-definidos e identidade Mackenzie."""

    def __init__(self):
        self._cor_primaria = _hex_para_rgb(settings.COR_PRIMARIA_MACKENZIE)
        self._cor_secundaria = _hex_para_rgb(settings.COR_SECUNDARIA_MACKENZIE)
        self._cor_fundo = _hex_para_rgb(settings.COR_FUNDO_MACKENZIE)
        self._cor_texto = _hex_para_rgb(settings.COR_TEXTO_MACKENZIE)
        self._instituicao = settings.NOME_INSTITUICAO

        self._handlers: Dict[TipoSlide, Callable] = {
            TipoSlide.CAPA: self._slide_capa,
            TipoSlide.SUMARIO: self._slide_sumario,
            TipoSlide.METRICAS: self._slide_metricas,
            TipoSlide.DIAGRAMA: self._slide_diagrama,
            TipoSlide.STATUS: self._slide_status,
            TipoSlide.PROXIMOS_PASSOS: self._slide_proximos_passos,
            TipoSlide.TEXTO: self._slide_texto,
            TipoSlide.ENCERRAMENTO: self._slide_encerramento,
        }

    def renderizar(self, apresentacao: Apresentacao) -> bytes:
        if not apresentacao.titulo or not apresentacao.titulo.strip():
            raise ApresentacaoInvalidaError("Apresentacao precisa de titulo.")
        if not apresentacao.slides:
            raise ApresentacaoInvalidaError("Apresentacao precisa de ao menos 1 slide.")

        try:
            prs = Presentation()
            prs.slide_width = _LARGURA_SLIDE
            prs.slide_height = _ALTURA_SLIDE

            for slide_dom in apresentacao.slides:
                handler = self._handlers.get(slide_dom.tipo)
                if handler is None:
                    raise ApresentacaoInvalidaError(
                        f"Tipo de slide nao suportado: {slide_dom.tipo}"
                    )
                handler(prs, slide_dom, apresentacao)

            buffer = io.BytesIO()
            prs.save(buffer)
            return buffer.getvalue()
        except ApresentacaoInvalidaError:
            raise
        except Exception as e:
            raise RenderizadorError(f"Falha ao renderizar PPTX: {e}") from e

    # ------------------------------------------------------------------
    # Layouts
    # ------------------------------------------------------------------

    def _slide_capa(self, prs, slide_dom: Slide, apresentacao: Apresentacao):
        slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_BRANCO])
        self._barra_lateral(slide)
        self._marca_dagua_rodape(slide)

        titulo = slide_dom.titulo or apresentacao.titulo
        subtitulo = slide_dom.subtitulo or apresentacao.subtitulo

        self._texto(
            slide, titulo,
            left=Inches(1.5), top=Inches(2.5), width=Inches(11), height=Inches(1.2),
            tamanho=44, negrito=True, cor=self._cor_primaria,
        )
        if subtitulo:
            self._texto(
                slide, subtitulo,
                left=Inches(1.5), top=Inches(3.7), width=Inches(11), height=Inches(0.8),
                tamanho=24, cor=self._cor_secundaria,
            )
        rodape = " | ".join(filter(None, [apresentacao.autor, apresentacao.data]))
        if rodape:
            self._texto(
                slide, rodape,
                left=Inches(1.5), top=Inches(5.5), width=Inches(11), height=Inches(0.5),
                tamanho=16, cor=self._cor_secundaria,
            )

    def _slide_sumario(self, prs, slide_dom: Slide, apresentacao: Apresentacao):
        slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_BRANCO])
        self._cabecalho(slide, slide_dom.titulo or "Sumario")
        self._marca_dagua_rodape(slide)

        itens = slide_dom.conteudo.get("itens") or []
        textos = [f"{i + 1}.  {item}" for i, item in enumerate(itens)]
        self._lista(slide, textos, top=Inches(1.6), tamanho=22)

    def _slide_metricas(self, prs, slide_dom: Slide, apresentacao: Apresentacao):
        slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_BRANCO])
        self._cabecalho(slide, slide_dom.titulo or "Metricas")
        self._marca_dagua_rodape(slide)

        itens = slide_dom.conteudo.get("itens") or []
        if not itens:
            return

        # Layout em grid 2 colunas
        cols = 2 if len(itens) > 1 else 1
        largura_cell = Inches(5.5)
        altura_cell = Inches(2.0)
        gap = Inches(0.4)
        origem_left = Inches(1.0)
        origem_top = Inches(1.8)

        for idx, item in enumerate(itens):
            row, col = divmod(idx, cols)
            left = origem_left + col * (largura_cell + gap)
            top = origem_top + row * (altura_cell + gap)
            self._cartao_metrica(
                slide,
                rotulo=str(item.get("rotulo", "")),
                valor=str(item.get("valor", "")),
                delta=str(item.get("delta", "")),
                left=left, top=top, width=largura_cell, height=altura_cell,
            )

    def _slide_diagrama(self, prs, slide_dom: Slide, apresentacao: Apresentacao):
        slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_BRANCO])
        self._cabecalho(slide, slide_dom.titulo or "Diagrama")
        self._marca_dagua_rodape(slide)

        imagem_b64 = slide_dom.conteudo.get("imagem_b64")
        if imagem_b64:
            try:
                blob = base64.b64decode(imagem_b64, validate=True)
                slide.shapes.add_picture(
                    io.BytesIO(blob),
                    left=Inches(2.0), top=Inches(1.6),
                    width=Inches(9.3), height=Inches(4.8),
                )
            except Exception:
                self._texto(
                    slide, "[imagem invalida — base64 nao decodificou]",
                    left=Inches(2.0), top=Inches(3.0), width=Inches(9.3), height=Inches(0.5),
                    tamanho=18, cor=self._cor_secundaria,
                )

        legenda = slide_dom.conteudo.get("legenda", "")
        if legenda:
            self._texto(
                slide, legenda,
                left=Inches(2.0), top=Inches(6.6), width=Inches(9.3), height=Inches(0.5),
                tamanho=14, italico=True, cor=self._cor_secundaria,
            )

    def _slide_status(self, prs, slide_dom: Slide, apresentacao: Apresentacao):
        slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_BRANCO])
        self._cabecalho(slide, slide_dom.titulo or "Status")
        self._marca_dagua_rodape(slide)

        itens = slide_dom.conteudo.get("itens") or []
        top = Inches(1.8)
        altura_linha = Inches(0.7)
        for item in itens:
            cor_indicador = _CORES_INDICADOR.get(item.get("cor", "verde"), _CORES_INDICADOR["verde"])
            bola = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(1.0), top + Inches(0.15), Inches(0.4), Inches(0.4)
            )
            bola.fill.solid()
            bola.fill.fore_color.rgb = cor_indicador
            bola.line.fill.background()
            self._texto(
                slide, str(item.get("texto", "")),
                left=Inches(1.6), top=top, width=Inches(11), height=altura_linha,
                tamanho=20, cor=self._cor_texto,
            )
            top += altura_linha

    def _slide_proximos_passos(self, prs, slide_dom: Slide, apresentacao: Apresentacao):
        slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_BRANCO])
        self._cabecalho(slide, slide_dom.titulo or "Proximos Passos")
        self._marca_dagua_rodape(slide)

        itens = slide_dom.conteudo.get("itens") or []
        textos = [f"{i + 1}.  {item}" for i, item in enumerate(itens)]
        self._lista(slide, textos, top=Inches(1.8), tamanho=22)

    def _slide_texto(self, prs, slide_dom: Slide, apresentacao: Apresentacao):
        slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_BRANCO])
        self._cabecalho(slide, slide_dom.titulo or "")
        self._marca_dagua_rodape(slide)

        paragrafos = slide_dom.conteudo.get("paragrafos") or []
        self._lista(slide, paragrafos, top=Inches(1.8), tamanho=18, marcador=False)

    def _slide_encerramento(self, prs, slide_dom: Slide, apresentacao: Apresentacao):
        slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_BRANCO])
        self._barra_lateral(slide)
        self._marca_dagua_rodape(slide)

        titulo = slide_dom.titulo or "Obrigado"
        self._texto(
            slide, titulo,
            left=Inches(1.5), top=Inches(3.0), width=Inches(11), height=Inches(1.5),
            tamanho=54, negrito=True, cor=self._cor_primaria,
        )
        if slide_dom.subtitulo:
            self._texto(
                slide, slide_dom.subtitulo,
                left=Inches(1.5), top=Inches(4.5), width=Inches(11), height=Inches(0.8),
                tamanho=22, cor=self._cor_secundaria,
            )

    # ------------------------------------------------------------------
    # Helpers de identidade visual
    # ------------------------------------------------------------------

    def _cabecalho(self, slide, titulo: str):
        barra = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), _LARGURA_SLIDE, Inches(0.6)
        )
        barra.fill.solid()
        barra.fill.fore_color.rgb = self._cor_primaria
        barra.line.fill.background()
        self._texto(
            slide, titulo,
            left=Inches(0.6), top=Inches(0.7), width=Inches(12), height=Inches(0.7),
            tamanho=28, negrito=True, cor=self._cor_primaria,
        )

    def _barra_lateral(self, slide):
        barra = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.6), _ALTURA_SLIDE
        )
        barra.fill.solid()
        barra.fill.fore_color.rgb = self._cor_primaria
        barra.line.fill.background()

    def _marca_dagua_rodape(self, slide):
        self._texto(
            slide, self._instituicao,
            left=Inches(0.5), top=Inches(7.0), width=Inches(12.3), height=Inches(0.4),
            tamanho=10, cor=self._cor_secundaria, alinhar_direita=True,
        )

    def _cartao_metrica(self, slide, rotulo: str, valor: str, delta: str, left, top, width, height):
        cartao = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        cartao.fill.solid()
        cartao.fill.fore_color.rgb = self._cor_fundo
        cartao.line.color.rgb = self._cor_primaria
        cartao.line.width = Pt(1.5)

        self._texto(
            slide, rotulo, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.5),
            tamanho=14, cor=self._cor_secundaria,
        )
        self._texto(
            slide, valor, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6), Inches(0.8),
            tamanho=36, negrito=True, cor=self._cor_primaria,
        )
        if delta:
            self._texto(
                slide, delta, left + Inches(0.3), top + Inches(1.5), width - Inches(0.6), Inches(0.4),
                tamanho=14, cor=self._cor_texto,
            )

    def _texto(self, slide, conteudo: str, left, top, width, height,
               tamanho=18, negrito=False, italico=False, cor=None, alinhar_direita=False):
        from pptx.enum.text import PP_ALIGN
        caixa = slide.shapes.add_textbox(left, top, width, height)
        tf = caixa.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if alinhar_direita:
            p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = conteudo
        run.font.size = Pt(tamanho)
        run.font.bold = negrito
        run.font.italic = italico
        if cor is not None:
            run.font.color.rgb = cor

    def _lista(self, slide, itens, top, tamanho=20, marcador=True):
        caixa = slide.shapes.add_textbox(Inches(1.0), top, Inches(11.3), Inches(5.0))
        tf = caixa.text_frame
        tf.word_wrap = True
        for i, texto in enumerate(itens):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = (f"•  {texto}" if marcador else str(texto))
            run.font.size = Pt(tamanho)
            run.font.color.rgb = self._cor_texto
            p.space_after = Pt(8)
