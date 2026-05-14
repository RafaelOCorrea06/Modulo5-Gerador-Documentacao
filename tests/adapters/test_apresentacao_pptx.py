# Testes do adaptador AdaptadorPythonPPTX (US GD-07).

import io

import pytest
from pptx import Presentation as AbrirPptx

from app.adapters.driven.renderizadores.adaptador_pythonpptx import AdaptadorPythonPPTX
from app.domain.entidades.apresentacao import Apresentacao
from app.domain.entidades.slide import Slide, TipoSlide
from app.domain.excecoes import ApresentacaoInvalidaError


def _construir_apresentacao_completa() -> Apresentacao:
    return Apresentacao(
        titulo="Status do Projeto",
        subtitulo="Sprint 5",
        autor="Time de Engenharia",
        data="Maio 2026",
        slides=[
            Slide(tipo=TipoSlide.CAPA),
            Slide(
                tipo=TipoSlide.SUMARIO,
                titulo="Agenda",
                conteudo={"itens": ["Visao geral", "Metricas", "Status", "Proximos passos"]},
            ),
            Slide(
                tipo=TipoSlide.METRICAS,
                titulo="Indicadores",
                conteudo={
                    "itens": [
                        {"rotulo": "Velocidade", "valor": "120%", "delta": "+5%"},
                        {"rotulo": "Cobertura", "valor": "78%", "delta": "+3%"},
                        {"rotulo": "Bugs abertos", "valor": "12", "delta": "-4"},
                    ]
                },
            ),
            Slide(
                tipo=TipoSlide.STATUS,
                titulo="Status dos modulos",
                conteudo={
                    "itens": [
                        {"texto": "Backend pronto", "cor": "verde"},
                        {"texto": "UI em revisao", "cor": "amarelo"},
                        {"texto": "Integracao bloqueada", "cor": "vermelho"},
                    ]
                },
            ),
            Slide(
                tipo=TipoSlide.PROXIMOS_PASSOS,
                titulo="Proximos passos",
                conteudo={"itens": ["Fechar IN-11", "Iniciar IA-10", "Demo dia 20"]},
            ),
            Slide(
                tipo=TipoSlide.TEXTO,
                titulo="Notas",
                conteudo={"paragrafos": ["Linha 1.", "Linha 2."]},
            ),
            Slide(tipo=TipoSlide.ENCERRAMENTO, titulo="Obrigado", subtitulo="Perguntas?"),
        ],
    )


def test_renderiza_pptx_valido_com_todos_os_layouts():
    bytes_pptx = AdaptadorPythonPPTX().renderizar(_construir_apresentacao_completa())

    # Magic bytes de zip (PPTX e um zip OOXML)
    assert bytes_pptx[:2] == b"PK"

    prs = AbrirPptx(io.BytesIO(bytes_pptx))
    assert len(prs.slides) == 7


def test_titulo_da_capa_aparece_no_slide():
    apres = _construir_apresentacao_completa()
    bytes_pptx = AdaptadorPythonPPTX().renderizar(apres)
    prs = AbrirPptx(io.BytesIO(bytes_pptx))

    capa = prs.slides[0]
    textos = [
        run.text
        for shape in capa.shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    ]
    assert any("Status do Projeto" in t for t in textos)


def test_apresentacao_sem_titulo_falha():
    apres = Apresentacao(titulo="", slides=[Slide(tipo=TipoSlide.CAPA)])
    with pytest.raises(ApresentacaoInvalidaError):
        AdaptadorPythonPPTX().renderizar(apres)


def test_apresentacao_sem_slides_falha():
    apres = Apresentacao(titulo="x", slides=[])
    with pytest.raises(ApresentacaoInvalidaError):
        AdaptadorPythonPPTX().renderizar(apres)


def test_diagrama_com_imagem_b64_invalida_nao_quebra():
    apres = Apresentacao(
        titulo="x",
        slides=[Slide(tipo=TipoSlide.DIAGRAMA, conteudo={"imagem_b64": "naoehbase64", "legenda": "fig 1"})],
    )
    # Nao deve levantar — apenas inserir mensagem de fallback no slide.
    bytes_pptx = AdaptadorPythonPPTX().renderizar(apres)
    assert bytes_pptx[:2] == b"PK"
